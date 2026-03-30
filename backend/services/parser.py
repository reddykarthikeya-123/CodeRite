"""Module for parsing various file formats and extracting text and images.

Supports PDF, DOCX, PPTX, Excel, and plain text files. Uses OCR for images
embedded in documents when necessary.
"""
from fastapi import UploadFile, HTTPException
import pdfplumber
from docx import Document
from pptx import Presentation
import io
import pandas as pd
from PIL import Image
import pytesseract
from pdf2image import convert_from_bytes
import os
import base64
import logging
import re
import asyncio
import shutil
import tempfile
import subprocess
import platform
import time
from pathlib import Path
from typing import Tuple, List, Dict, Any
import xml.etree.ElementTree as ET
import zipfile
from config.logging_config import get_logger

logger = get_logger(__name__)

# Try to import python-magic for MIME type detection
# This is optional - if not available, we'll skip MIME validation
try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False
    logger.warning("python-magic not available. MIME type validation will be skipped. Install libmagic if needed.")

# Max file size 50MB
MAX_FILE_SIZE = 50 * 1024 * 1024

# Whitelisted extensions
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".py", ".js", ".ts", ".json", ".html", ".css", ".xlsx", ".xls", ".csv", ".car"}

# Allowed MIME types
ALLOWED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "text/plain",
    "text/markdown",
    "text/x-python",
    "application/javascript",
    "application/json",
    "text/html",
    "text/css",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
    "text/csv",
    "application/zip",
    "application/octet-stream",
    "application/x-zip-compressed"
}


def _is_truthy_env(value: str) -> bool:
    return value.strip().lower() in {"1", "true", "yes", "on"}


def _docx_pagination_required() -> bool:
    return _is_truthy_env(os.getenv("DOCX_PAGINATION_REQUIRED", "false"))


def _default_pagination_metadata() -> Dict[str, Any]:
    return {
        "enabled": False,
        "format": None,
        "total_pages": None,
        "provider": "none",
        "warning": None
    }


def _build_pagination_metadata(
    enabled: bool,
    page_format: str | None,
    total_pages: int | None,
    provider: str,
    warning: str | None = None
) -> Dict[str, Any]:
    return {
        "enabled": enabled,
        "format": page_format,
        "total_pages": total_pages,
        "provider": provider,
        "warning": warning
    }


def _extract_total_pages(text: str) -> int:
    matches = re.findall(r'--- Page (\d+) (?:Text|Tables) ---', text)
    if not matches:
        return 0
    return max(int(page) for page in matches)


def _safe_int_env(name: str, default: int) -> int:
    raw = os.getenv(name, str(default))
    try:
        return int(raw)
    except (TypeError, ValueError):
        return default


def _safe_str(value: Any) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _extract_docx_custom_properties(content: bytes) -> Dict[str, str]:
    """Extract custom document properties from a DOCX package."""
    properties: Dict[str, str] = {}
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            if "docProps/custom.xml" not in archive.namelist():
                return properties

            xml_bytes = archive.read("docProps/custom.xml")
            root = ET.fromstring(xml_bytes)
            namespaces = {
                "cp": "http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"
            }

            for prop in root.findall("cp:property", namespaces):
                name = _safe_str(prop.attrib.get("name"))
                if not name:
                    continue
                value = ""
                for child in list(prop):
                    if child.text and child.text.strip():
                        value = child.text.strip()
                        break
                properties[name] = value
    except Exception as exc:
        logger.warning(f"Unable to read DOCX custom properties: {exc}")
    return properties


def _normalize_image_for_model(image: Image.Image) -> Image.Image:
    """Normalize images before OCR/model transport to keep payloads stable."""
    normalized = image
    if normalized.mode in ("RGBA", "P"):
        normalized = normalized.convert("RGB")
    elif normalized.mode != "RGB":
        normalized = normalized.convert("RGB")

    max_dimension = max(256, _safe_int_env("VISION_IMAGE_MAX_DIM", 1600))
    width, height = normalized.size
    longest_side = max(width, height)
    if longest_side > max_dimension:
        scale = max_dimension / float(longest_side)
        resized = (
            max(1, int(width * scale)),
            max(1, int(height * scale))
        )
        normalized = normalized.resize(resized, Image.LANCZOS)

    return normalized


def _prepare_image_for_model(image: Image.Image) -> tuple[Image.Image, str, int]:
    """Return normalized PIL image plus encoded JPEG payload details."""
    normalized = _normalize_image_for_model(image)
    jpeg_quality = max(40, min(95, _safe_int_env("VISION_IMAGE_JPEG_QUALITY", 80)))
    buffered = io.BytesIO()
    normalized.save(
        buffered,
        format="JPEG",
        quality=jpeg_quality,
        optimize=True
    )
    image_bytes = buffered.getvalue()
    return normalized, base64.b64encode(image_bytes).decode("utf-8"), len(image_bytes)


def _get_pdf_visual_counts(page: pdfplumber.page.Page) -> Dict[str, int]:
    """Collect page-level visual object counts for OCR gating and LLM grounding."""
    return {
        "image_objects": len(page.images),
        "line_objects": len(page.lines),
        "rect_objects": len(page.rects),
        "curve_objects": len(page.curves),
    }


def _format_pdf_visual_metadata(page_number: int, visual_counts: Dict[str, int]) -> str:
    """Summarize page-level visual object counts for LLM grounding."""
    image_objects = visual_counts["image_objects"]
    line_objects = visual_counts["line_objects"]
    rect_objects = visual_counts["rect_objects"]
    curve_objects = visual_counts["curve_objects"]

    return (
        f"\n--- Page {page_number} Visual Metadata ---\n"
        f"image_objects={image_objects}\n"
        f"line_objects={line_objects}\n"
        f"rect_objects={rect_objects}\n"
        f"curve_objects={curve_objects}\n"
    )


def _resolve_soffice_path() -> str:
    configured_path = os.getenv("SOFFICE_PATH", "").strip()
    if configured_path:
        if os.path.exists(configured_path):
            return configured_path
        raise RuntimeError(f"SOFFICE_PATH is set but not found: {configured_path}")

    discovered_path = shutil.which("soffice")
    if discovered_path:
        return discovered_path

    if platform.system() == "Windows":
        candidate_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        for candidate in candidate_paths:
            if os.path.exists(candidate):
                return candidate

    raise RuntimeError(
        "LibreOffice 'soffice' binary not found. Install LibreOffice and set SOFFICE_PATH if needed."
    )


async def _convert_docx_to_pdf_with_libreoffice(content: bytes) -> bytes:
    """Converts DOCX bytes to PDF bytes using headless LibreOffice."""
    soffice_path = _resolve_soffice_path()
    timeout_sec = int(os.getenv("DOCX_CONVERT_TIMEOUT_SEC", "90"))
    start_time = time.perf_counter()
    logger.info(
        f"DOCX->PDF conversion started via LibreOffice. soffice='{soffice_path}', timeout={timeout_sec}s"
    )

    with tempfile.TemporaryDirectory(prefix="docx_convert_") as tmpdir:
        input_path = os.path.join(tmpdir, "input.docx")
        output_dir = os.path.join(tmpdir, "out")
        profile_dir = os.path.join(tmpdir, "profile")

        os.makedirs(output_dir, exist_ok=True)
        os.makedirs(profile_dir, exist_ok=True)

        with open(input_path, "wb") as f:
            f.write(content)

        # Isolate LO profile per request to avoid lock/contention issues.
        profile_uri = Path(profile_dir).resolve().as_uri()
        command = [
            soffice_path,
            "--headless",
            "--nologo",
            "--nodefault",
            "--nolockcheck",
            "--norestore",
            f"-env:UserInstallation={profile_uri}",
            "--convert-to",
            "pdf:writer_pdf_Export",
            "--outdir",
            output_dir,
            input_path,
        ]

        completed = await asyncio.to_thread(
            subprocess.run,
            command,
            capture_output=True,
            text=True,
            check=False,
            timeout=timeout_sec,
        )

        if completed.returncode != 0:
            stderr = (completed.stderr or completed.stdout or "Unknown conversion error").strip()
            logger.error(f"DOCX->PDF conversion failed. LibreOffice stderr/stdout: {stderr}")
            raise RuntimeError(f"LibreOffice conversion failed with code {completed.returncode}: {stderr}")

        output_pdf_path = os.path.join(output_dir, "input.pdf")
        if not os.path.exists(output_pdf_path):
            stderr = (completed.stderr or completed.stdout or "No output PDF generated").strip()
            logger.error(f"DOCX->PDF conversion produced no output PDF. LibreOffice stderr/stdout: {stderr}")
            raise RuntimeError(f"LibreOffice conversion produced no PDF output: {stderr}")

        with open(output_pdf_path, "rb") as pdf_file:
            pdf_bytes = pdf_file.read()

        elapsed_ms = int((time.perf_counter() - start_time) * 1000)
        logger.info(
            f"DOCX->PDF conversion succeeded via LibreOffice in {elapsed_ms}ms. "
            f"output_size={len(pdf_bytes)} bytes"
        )
        return pdf_bytes

def sanitize_filename(filename: str) -> str:
    """Sanitize the filename by taking only the base name.

    Args:
        filename: The original filename from UploadFile.

    Returns:
        The sanitized base filename.
    """
    return os.path.basename(filename)

async def parse_file(file: UploadFile) -> dict:
    """Parse uploaded file and extract text and images with security validations.

    Args:
        file: The uploaded file object.

    Returns:
        A dictionary containing extracted 'text' and 'images'.

    Raises:
        HTTPException: For invalid file types, sizes, or parsing errors.
    """
    # 1. Filename sanitization
    filename = sanitize_filename(file.filename).lower()

    # 2. Extension validation
    _, ext = os.path.splitext(filename)
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"Unsupported file extension: {ext}")

    # 3. Size validation
    content = await file.read()
    if len(content) == 0:
        raise HTTPException(status_code=400, detail="Empty file uploaded")
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail="File size exceeds 50MB limit")

    # 4. MIME type validation using python-magic (MANDATORY for security)
    if not MAGIC_AVAILABLE:
        logger.error("python-magic is not available. This is a security requirement.")
        raise HTTPException(
            status_code=500,
            detail="Server configuration error: MIME validation unavailable. Please contact support."
        )
    
    mime_type = magic.from_buffer(content, mime=True)
    
    # Validate MIME type matches expected types
    if mime_type not in ALLOWED_MIME_TYPES:
        # Allow text/plain for code files
        if not (ext in [".py", ".js", ".ts", ".json", ".html", ".css", ".md", ".txt"] and mime_type.startswith("text/")):
            logger.warning(f"MIME type mismatch for '{filename}': detected '{mime_type}', expected one of {ALLOWED_MIME_TYPES}")
            raise HTTPException(
                status_code=400,
                detail=f"File type mismatch. The uploaded file appears to be a {mime_type}, not a valid {ext} file."
            )
    
    # 5. Additional content validation for PDFs
    if filename.endswith(".pdf"):
        if not content.startswith(b"%PDF"):
            raise HTTPException(status_code=400, detail="Invalid PDF file signature detected.")

    # Reset file pointer or use content directly for parsing
    # Since we already read it into 'content', we'll modify parsers to accept bytes or use io.BytesIO

    images = []
    text_content = ""
    pagination_metadata = _default_pagination_metadata()

    try:
        if filename.endswith(".pdf"):
            text_content, images = await _parse_pdf_from_bytes(content)
            total_pages = _extract_total_pages(text_content)
            pagination_metadata = _build_pagination_metadata(
                enabled=total_pages > 0,
                page_format="Page" if total_pages > 0 else None,
                total_pages=total_pages if total_pages > 0 else None,
                provider="native_pdf",
                warning=None
            )
        elif filename.endswith(".docx"):
            text_content, images, pagination_metadata = await _parse_docx_from_bytes(content)
            logger.info(
                "DOCX parse completed. "
                f"pagination_enabled={pagination_metadata.get('enabled')} "
                f"provider={pagination_metadata.get('provider')} "
                f"total_pages={pagination_metadata.get('total_pages')} "
                f"warning={pagination_metadata.get('warning')}"
            )
        elif filename.endswith(".pptx"):
            text_content, images = await _parse_pptx_from_bytes(content)
        elif filename.endswith((".txt", ".md", ".py", ".js", ".ts", ".json", ".html", ".css")):
            text_content = content.decode("utf-8")
            images = []
        elif filename.endswith((".xlsx", ".xls", ".csv")):
            text_content = await _parse_excel_from_bytes(content, filename)
            images = []
        elif filename.endswith(".car"):
            # Returns structured data for better chunking
            parsed_data = await _parse_car_from_bytes(content)
            # Convert to text format for backward compatibility
            text_parts = []
            for file_info in parsed_data["files"]:
                text_parts.append(f"\n--- File: {file_info['filename']} ---\n{file_info['content']}")
            text_content = "\n".join(text_parts)
            images = parsed_data["images"]
            # Store structured data for AI engine to use
            text_content = f"{text_content}\n\n[CAR_METADATA] total_size={parsed_data['total_size']}, file_count={parsed_data['file_count']} [/CAR_METADATA]"

        logger.info(
            "File parse completed: "
            f"filename={filename} "
            f"text_chars={len(text_content)} "
            f"images={len(images)} "
            f"image_chars_total={sum(len(image) for image in images)} "
            f"pagination_enabled={pagination_metadata.get('enabled')} "
            f"pagination_provider={pagination_metadata.get('provider')}"
        )
        return {"text": text_content, "images": images, "pagination_metadata": pagination_metadata}
    except Exception as e:
        logger.error(f"Error parsing file '{filename}': {e}", exc_info=True)
        raise HTTPException(status_code=500, detail="Error parsing file. Please check the file format and try again.")

# Helper functions to handle bytes instead of UploadFile to avoid double reading
async def _parse_pdf_from_bytes(content: bytes) -> Tuple[str, List[str]]:
    """Extracts text and images from a PDF provided as bytes.

    Args:
        content: The raw PDF file content.

    Returns:
        A tuple of (extracted_text, list_of_base64_images).
    """
    text_parts: List[str] = []
    base64_images: List[str] = []
    page_text_lengths: List[int] = []
    page_visual_counts: List[Dict[str, int]] = []
    table_rows = 0
    ocr_blocks = 0
    image_bytes_total = 0

    ocr_mode = os.getenv("PDF_OCR_MODE", "always").strip().lower()
    if ocr_mode not in {"always", "auto", "off"}:
        ocr_mode = "always"
    ocr_min_text_chars = _safe_int_env("PDF_OCR_MIN_TEXT_CHARS_PER_PAGE", 50)
    ocr_max_pages = _safe_int_env("PDF_OCR_MAX_PAGES", 100)
    ocr_visual_object_threshold = max(1, _safe_int_env("PDF_OCR_VISUAL_OBJECT_THRESHOLD", 8))
    render_dpi = max(72, _safe_int_env("PDF_RENDER_DPI", 160))

    try:
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text(layout=True) or ""
                page_text_clean = page_text.strip()
                page_text_lengths.append(len(page_text_clean))
                visual_counts = _get_pdf_visual_counts(page)
                page_visual_counts.append(visual_counts)
                if page_text_clean:
                    text_parts.append(f"\n--- Page {i+1} Text ---\n{page_text}\n")

                tables = page.extract_tables()
                if tables:
                    text_parts.append(f"\n--- Page {i+1} Tables ---\n")
                    for table_idx, table in enumerate(tables):
                        text_parts.append(f"Table {table_idx + 1}:\n")
                        for row in table:
                            cleaned_row = [str(cell).replace("\n", " ").strip() if cell is not None else "" for cell in row]
                            text_parts.append("| " + " | ".join(cleaned_row) + " |\n")
                            table_rows += 1
                        text_parts.append("\n")

                text_parts.append(_format_pdf_visual_metadata(i + 1, visual_counts))

        # Process images with OCR in non-blocking manner
        images = convert_from_bytes(content, dpi=render_dpi)
        for i, img in enumerate(images):
            normalized_img, img_b64, image_bytes = _prepare_image_for_model(img)
            base64_images.append(img_b64)
            image_bytes_total += image_bytes

            if ocr_mode == "off":
                should_ocr = False
            elif ocr_mode == "always":
                should_ocr = True
            else:
                page_text_len = page_text_lengths[i] if i < len(page_text_lengths) else 0
                visual_counts = page_visual_counts[i] if i < len(page_visual_counts) else {}
                visual_object_count = sum(int(value) for value in visual_counts.values())
                has_explicit_visual_artifact = (
                    int(visual_counts.get("image_objects", 0)) > 0 or
                    int(visual_counts.get("rect_objects", 0)) > 0 or
                    int(visual_counts.get("curve_objects", 0)) > 0
                )
                should_ocr = (
                    page_text_len < ocr_min_text_chars or
                    has_explicit_visual_artifact or
                    visual_object_count >= ocr_visual_object_threshold
                )

            if should_ocr and ocr_blocks < ocr_max_pages:
                ocr_text = await asyncio.to_thread(pytesseract.image_to_string, img)
                if ocr_text.strip():
                    text_parts.append(f"\n--- Page {i+1} OCR ---\n{ocr_text}\n")
                    ocr_blocks += 1
    except Exception as e:
        logger.warning(f"OCR/Vision warning on PDF: {e}")

    text = "".join(text_parts)
    logger.info(
        "PDF parse stats: "
        f"pages={len(page_text_lengths)} "
        f"tables_rows={table_rows} "
        f"images_extracted={len(base64_images)} "
        f"image_bytes_total={image_bytes_total} "
        f"ocr_blocks={ocr_blocks} "
        f"ocr_mode={ocr_mode} "
        f"ocr_visual_object_threshold={ocr_visual_object_threshold} "
        f"render_dpi={render_dpi} "
        f"text_chars={len(text)}"
    )
    return text, base64_images


async def _parse_docx_from_bytes(content: bytes) -> Tuple[str, List[str], Dict[str, Any]]:
    """Extracts text and images from a DOCX provided as bytes.

    Args:
        content: The raw DOCX file content.

    Returns:
        A tuple of (extracted_text, list_of_base64_images, pagination_metadata).
    """
    conversion_error: str | None = None

    # Preferred path: convert DOCX to PDF and reuse PDF parser for page-accurate markers.
    logger.info("DOCX pagination: attempting LibreOffice conversion for page-accurate references.")
    try:
        converted_pdf_bytes = await _convert_docx_to_pdf_with_libreoffice(content)
        text, base64_images = await _parse_pdf_from_bytes(converted_pdf_bytes)
        total_pages = _extract_total_pages(text)
        if total_pages <= 0:
            conversion_error = "Converted PDF did not contain usable page markers."
            raise RuntimeError(conversion_error)

        pagination_metadata = _build_pagination_metadata(
            enabled=True,
            page_format="Page",
            total_pages=total_pages,
            provider="libreoffice_pdf",
            warning=None,
        )
        logger.info(
            f"DOCX pagination enabled via LibreOffice PDF conversion. total_pages={total_pages}, "
            f"images_extracted={len(base64_images)}"
        )
        return text, base64_images, pagination_metadata
    except Exception as conversion_exception:
        conversion_error = str(conversion_exception)
        logger.warning(f"DOCX pagination conversion failed. Falling back to text extraction: {conversion_error}")
        if _docx_pagination_required():
            raise RuntimeError(
                f"DOCX pagination is required but DOCX->PDF conversion failed: {conversion_error}"
            ) from conversion_exception

    # Fallback path: extract text directly without page references.
    doc = Document(io.BytesIO(content))
    text_parts: List[str] = []
    paragraph_count = 0
    table_row_count = 0
    header_paragraph_count = 0
    footer_paragraph_count = 0
    ocr_blocks = 0
    image_bytes_total = 0

    # Include document properties for better metadata detection (version/author/etc.).
    core = doc.core_properties
    core_properties = {
        "title": _safe_str(core.title),
        "subject": _safe_str(core.subject),
        "author": _safe_str(core.author),
        "last_modified_by": _safe_str(core.last_modified_by),
        "revision": _safe_str(core.revision),
        "created": _safe_str(core.created),
        "modified": _safe_str(core.modified),
        "category": _safe_str(core.category),
        "keywords": _safe_str(core.keywords),
    }
    custom_properties = _extract_docx_custom_properties(content)

    text_parts.append("\n--- DOCX Core Properties ---\n")
    for key, value in core_properties.items():
        if value:
            text_parts.append(f"{key}: {value}\n")

    if custom_properties:
        text_parts.append("\n--- DOCX Custom Properties ---\n")
        for key, value in custom_properties.items():
            text_parts.append(f"{key}: {value}\n")

    text_parts.append("\n--- DOCX Body Paragraphs ---\n")
    for para_index, para in enumerate(doc.paragraphs, start=1):
        para_text = para.text.strip()
        if para_text:
            paragraph_count += 1
            text_parts.append(f"P{para_index}: {para_text}\n")

    for table_index, table in enumerate(doc.tables, start=1):
        text_parts.append(f"\n--- DOCX Table {table_index} ---\n")
        for row_index, row in enumerate(table.rows, start=1):
            cleaned_row = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            if any(cleaned_row):
                table_row_count += 1
                text_parts.append(f"Row {row_index}: | " + " | ".join(cleaned_row) + " |\n")

    for section_index, section in enumerate(doc.sections, start=1):
        header_lines = [p.text.strip() for p in section.header.paragraphs if p.text and p.text.strip()]
        if header_lines:
            text_parts.append(f"\n--- DOCX Header Section {section_index} ---\n")
            for line in header_lines:
                header_paragraph_count += 1
                text_parts.append(line + "\n")

        footer_lines = [p.text.strip() for p in section.footer.paragraphs if p.text and p.text.strip()]
        if footer_lines:
            text_parts.append(f"\n--- DOCX Footer Section {section_index} ---\n")
            for line in footer_lines:
                footer_paragraph_count += 1
                text_parts.append(line + "\n")

    base64_images: List[str] = []
    try:
        for image_index, rel in enumerate(doc.part.rels.values(), start=1):
            if "image" in rel.target_ref:
                img_data = rel.target_part.blob
                img = Image.open(io.BytesIO(img_data))
                normalized_img, img_b64, image_bytes = _prepare_image_for_model(img)
                base64_images.append(img_b64)
                image_bytes_total += image_bytes

                ocr_text = await asyncio.to_thread(pytesseract.image_to_string, normalized_img)
                if ocr_text.strip():
                    ocr_blocks += 1
                    text_parts.append(f"\n--- DOCX Embedded Image OCR {image_index} ---\n{ocr_text}\n")
    except Exception as e:
        logger.warning(f"OCR/Vision warning on Docx: {e}")

    text = "".join(text_parts)
    warning = (
        "Page references disabled for this file because Word-to-PDF pagination failed."
        if conversion_error else
        "Page references disabled for this file."
    )
    pagination_metadata = _build_pagination_metadata(
        enabled=False,
        page_format=None,
        total_pages=None,
        provider="none",
        warning=warning
    )
    logger.info(
        "DOCX pagination disabled; using fallback text extraction without page references. "
        f"paragraphs={paragraph_count} "
        f"table_rows={table_row_count} "
        f"header_paragraphs={header_paragraph_count} "
        f"footer_paragraphs={footer_paragraph_count} "
        f"core_properties={sum(1 for value in core_properties.values() if value)} "
        f"custom_properties={len(custom_properties)} "
        f"fallback_images_extracted={len(base64_images)} "
        f"image_bytes_total={image_bytes_total} "
        f"ocr_blocks={ocr_blocks} "
        f"text_chars={len(text)}"
    )
    return text, base64_images, pagination_metadata


async def _parse_pptx_from_bytes(content: bytes) -> Tuple[str, List[str]]:
    """Extracts text and images from a PPTX provided as bytes.

    Args:
        content: The raw PPTX file content.

    Returns:
        A tuple of (extracted_text, list_of_base64_images).
    """
    prs = Presentation(io.BytesIO(content))
    text = ""
    base64_images = []
    image_bytes_total = 0
    ocr_blocks = 0

    for i, slide in enumerate(prs.slides):
        text += f"\n--- Slide {i+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
            if getattr(shape, "shape_type", None) == 13: 
                try:
                    img_bytes = shape.image.blob
                    img = Image.open(io.BytesIO(img_bytes))
                    normalized_img, img_b64, image_bytes = _prepare_image_for_model(img)
                    base64_images.append(img_b64)
                    image_bytes_total += image_bytes
                    ocr_text = await asyncio.to_thread(pytesseract.image_to_string, normalized_img)
                    if ocr_text.strip():
                        ocr_blocks += 1
                        text += f"\n[Embedded Image OCR]: {ocr_text}\n"
                except Exception as e:
                    logger.error(f"Failed to process image on slide {i+1}: {e}")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                text += f"\n[Speaker Notes]:\n{notes}\n"

    logger.info(
        "PPTX parse stats: "
        f"slides={len(prs.slides)} "
        f"images_extracted={len(base64_images)} "
        f"image_bytes_total={image_bytes_total} "
        f"ocr_blocks={ocr_blocks} "
        f"text_chars={len(text)}"
    )
    return text, base64_images

async def _parse_excel_from_bytes(content: bytes, filename: str) -> str:
    """Extracts text from an Excel or CSV file provided as bytes.

    Args:
        content: The raw file content.
        filename: The filename to determine if it's CSV or Excel.

    Returns:
        The extracted text content formatted as CSV.
    """
    text = ""
    if filename.endswith(".csv"):
        df = pd.read_csv(io.BytesIO(content))
        text += f"--- CSV Data ---\n{df.to_csv(index=False)}\n"
    else:
        xls = pd.ExcelFile(io.BytesIO(content))
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            text += f"\n--- Excel Sheet: {sheet_name} ---\n"
            text += df.to_csv(index=False) + "\n"
    return text

async def _parse_car_from_bytes(content: bytes) -> dict:
    """Recursively extract XML, XSL, WSDL, and properties from .car and .iar archives.
    
    Returns structured data with individual files preserved for better chunking.
    """
    files = []
    images = []
    total_size = 0

    def process_zip_content(zip_bytes, prefix=""):
        nonlocal total_size
        try:
            with zipfile.ZipFile(io.BytesIO(zip_bytes)) as z:
                for info in z.infolist():
                    if info.is_dir():
                        continue
                    filename = info.filename.lower()

                    if filename.endswith(('.xml', '.xsl', '.wsdl', '.properties', '.jpr', '.jca', '.xqy')):
                        file_data = z.read(info.filename)
                        try:
                            decoded = file_data.decode('utf-8')
                        except UnicodeDecodeError:
                            decoded = file_data.decode('latin-1', errors='ignore')

                        # Store each file separately with its path
                        files.append({
                            "filename": f"{prefix}{info.filename}",
                            "content": decoded
                        })
                        total_size += len(decoded)

                    elif filename.endswith('.iar'):
                        file_data = z.read(info.filename)
                        process_zip_content(file_data, prefix + info.filename + " -> ")

        except Exception as e:
            logger.error(f"Error extracting archive: {str(e)}")
            files.append({
                "filename": "ERROR",
                "content": f"[Error extracting archive: {str(e)}]"
            })

    process_zip_content(content)
    
    # Add metadata
    return {
        "files": files,
        "total_size": total_size,
        "file_count": len(files),
        "images": images
    }
